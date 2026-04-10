# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import base64
from unittest.mock import MagicMock, patch

from pptx.util import Pt

from presentation_agent.shared_libraries.utils import (
    _apply_font,
    _decode_base64_to_bytes,
    _insert_image,
)


def test_decode_base64():
    test_str = "Hello World"
    b64_str = base64.b64encode(test_str.encode()).decode()
    decoded = _decode_base64_to_bytes(b64_str)
    assert decoded == test_str.encode()

    data_uri = f"data:image/png;base64,{b64_str}"
    decoded_uri = _decode_base64_to_bytes(data_uri)
    assert decoded_uri == test_str.encode()

    assert _decode_base64_to_bytes("invalid_base64!!!") is None


def test_apply_font():
    run = MagicMock()

    _apply_font(None, "Arial", 12.0, (255, 0, 0))

    _apply_font(run, "Arial", 12.0, (255, 0, 0))
    assert run.font.name == "Arial"
    assert run.font.size == Pt(12.0)

    _apply_font(run, "Arial", 12.0, (255, 0))

    # Trigger exception
    with patch(
        "presentation_agent.shared_libraries.utils.Pt",
        side_effect=Exception("Test Error"),
    ):
        _apply_font(run, "Arial", 12.0, (255, 0, 0))


@patch("presentation_agent.shared_libraries.utils.get_gcs_client")
@patch("presentation_agent.shared_libraries.utils.Image.open")
@patch("os.path.exists")
@patch("builtins.open", new_callable=MagicMock)
def test_insert_image(
    mock_open, mock_exists, mock_image_open, mock_get_gcs_client
):
    prs = MagicMock()
    prs.slide_width = 9600000
    prs.slide_height = 5400000
    slide = MagicMock()

    mock_img = MagicMock()
    mock_img.size = (100, 100)
    mock_image_open.return_value.__enter__.return_value = mock_img

    _insert_image(prs, slide, b"fake_bytes", (100, 100, 200, 200))

    b64_str = base64.b64encode(b"fake_bytes").decode()
    _insert_image(prs, slide, b64_str, (100, 100, 200, 200))

    mock_client = MagicMock()
    mock_blob = MagicMock()
    mock_blob.download_as_bytes.return_value = b"gcs_bytes"
    mock_client.bucket.return_value.blob.return_value = mock_blob
    mock_get_gcs_client.return_value = mock_client

    _insert_image(prs, slide, "gs://bucket/blob.png", (100, 100, 200, 200))

    # Exception inside GCS download
    mock_blob.download_as_bytes.side_effect = Exception("Download error")
    _insert_image(prs, slide, "gs://bucket/blob.png", (100, 100, 200, 200))

    mock_get_gcs_client.return_value = None
    _insert_image(prs, slide, "gs://bucket/blob.png", (100, 100, 200, 200))

    mock_exists.return_value = True
    mock_open.return_value.__enter__.return_value.read.return_value = (
        b"local_bytes"
    )
    _insert_image(prs, slide, "/path/to/image.png", (100, 100, 200, 200))

    mock_exists.return_value = False
    _insert_image(prs, slide, "/path/to/missing.png", (100, 100, 200, 200))

    _insert_image(prs, slide, 12345, (100, 100, 200, 200))

    slide.shapes.add_picture.reset_mock()
    _insert_image(prs, slide, b"fake_bytes", None)

    mock_img.size = (0, 0)
    slide.shapes.add_picture.reset_mock()
    _insert_image(prs, slide, b"fake_bytes", None)

    # Test empty image bytes (local file is empty)
    mock_exists.return_value = True
    mock_open.return_value.__enter__.return_value.read.return_value = b""
    _insert_image(prs, slide, "/path/to/empty.png", (100, 100, 200, 200))
