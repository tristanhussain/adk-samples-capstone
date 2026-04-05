# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

from pptx import Presentation

from presentation_agent.tools import get_safe_layout


def test_get_safe_layout_exact_match():
    prs = Presentation()
    # Usually index 0 is "Title Slide"
    layout = get_safe_layout(prs, "Title Slide")
    assert layout.name == "Title Slide"


def test_get_safe_layout_partial_match():
    prs = Presentation()
    # Even if we misspell or hallucinate slightly, it should find it
    layout = get_safe_layout(prs, "Some Title layout")
    assert layout.name == "Title Slide"


def test_get_safe_layout_fallback():
    prs = Presentation()
    # A completely bogus name should fallback to index 1 (usually "Title and Content")
    layout = get_safe_layout(prs, "NonExistentLayout")
    assert layout.name == "Title and Content"


def test_get_safe_layout_empty():
    prs = Presentation()
    layout = get_safe_layout(prs, "")
    assert layout.name == "Title Slide"
