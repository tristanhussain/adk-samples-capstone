# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import os
import tempfile

from pptx import Presentation

from presentation_agent.tools.pptx_editor import (
    add_slide_to_end,
    delete_slide,
    edit_slide_text,
    read_presentation_details,
    read_presentation_outline,
)


def create_sample_presentation():
    """Helper to create a simple presentation for testing."""
    prs = Presentation()
    # Add Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sample Title"
    subtitle.text = "Sample Subtitle"

    # Add Content Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Slide 2 Title"
    content.text = "Old Text"

    # Save to temp file
    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path


def test_edit_slide_text():
    pptx_path = create_sample_presentation()
    try:
        # The function signature is (pptx_path, slide_number, new_title, new_bullets)
        result = edit_slide_text(
            pptx_path=pptx_path,
            slide_number=2,
            new_title="New Title",
            new_bullets=["New Bullet 1", "New Bullet 2"],
        )

        assert "Successfully edited slide" in result

        # Verify changes
        prs = Presentation(pptx_path)
        slide = prs.slides[1]  # 0-indexed internally
        found_title = False
        found_bullet = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "New Title" in shape.text:
                    found_title = True
                if "New Bullet 1" in shape.text:
                    found_bullet = True

        assert found_title, "The new title was not found on the slide."
        assert found_bullet, "The new bullets were not found on the slide."
    finally:
        os.remove(pptx_path)


def test_add_slide_to_end():
    pptx_path = create_sample_presentation()
    try:
        # Before adding, we have 2 slides
        prs = Presentation(pptx_path)
        assert len(prs.slides) == 2

        result = add_slide_to_end(pptx_path, "Title and Content")
        assert "Successfully added new slide" in result

        # After adding, we should have 3 slides
        prs = Presentation(pptx_path)
        assert len(prs.slides) == 3
    finally:
        os.remove(pptx_path)


def test_delete_slide():
    pptx_path = create_sample_presentation()
    try:
        # Before deleting, we have 2 slides
        prs = Presentation(pptx_path)
        assert len(prs.slides) == 2

        result = delete_slide(pptx_path, 2)
        assert "Successfully deleted slide 2" in result

        # After deleting, we should have 1 slide
        prs = Presentation(pptx_path)
        assert len(prs.slides) == 1
        assert "Sample Title" in prs.slides[0].shapes.title.text
    finally:
        os.remove(pptx_path)


def test_read_presentation_outline():
    pptx_path = create_sample_presentation()
    try:
        outline = read_presentation_outline(pptx_path)
        assert "Slide 1:" in outline
        assert "Sample Title" in outline
        assert "Slide 2:" in outline
        assert "Slide 2 Title" in outline
    finally:
        os.remove(pptx_path)


def test_read_presentation_details():
    pptx_path = create_sample_presentation()
    try:
        details = read_presentation_details(pptx_path)
        assert "Slide 1 Title" in details
        assert "Sample Title" in details
        assert "Sample Subtitle" in details
        assert "Slide 2 Title" in details
        assert "Old Text" in details
    finally:
        os.remove(pptx_path)
