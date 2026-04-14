# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import base64
import io
import os
import re
from typing import Any

from PIL import Image

# Presentation & Image Imports
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Local Application Imports
from .config import get_gcs_client, get_logger


def _apply_font(
    run,
    name: str | None,
    size_pt: float | None,
    color_rgb: tuple[int, int, int] | None,
):
    """Safely applies font style (name, size, color) to a given text run."""
    log = get_logger("apply_font")
    if not run:
        log.warning("Attempted to apply font to a non-existent run.")
        return
    try:
        if name:
            run.font.name = name
        if size_pt:
            run.font.size = Pt(size_pt)
        if color_rgb:
            if isinstance(color_rgb, (list, tuple)) and len(color_rgb) == 3:
                run.font.color.rgb = RGBColor(*color_rgb)
            else:
                log.warning(
                    f"Invalid color_rgb format: {color_rgb}. Expected (R, G, B)."
                )
    except Exception as e:
        log.error(f"Error applying font style: {e}", exc_info=True)


def _decode_base64_to_bytes(base64_string: str) -> bytes | None:
    """Decodes a base64 string to bytes, returning None if invalid."""
    try:
        if base64_string.startswith("data:image/"):
            base64_string = base64_string.split(",", 1)[1]
        return base64.b64decode(base64_string, validate=True)
    except (base64.binascii.Error, ValueError) as e:
        get_logger("decode_base64").error(
            f"Failed to decode base64 string: {e}"
        )
        return None


def _insert_image(
    prs: Presentation,
    slide: Any,
    image_data: str | bytes,
    box_hint: tuple[int, int, int, int] | None,
):
    """
    Loads an image from various sources and inserts it into a slide.
    If a box_hint is provided, it scales the image to fit within the box
    while preserving its aspect ratio, preventing distortion.
    If no box_hint is provided, it centers the image on the slide.
    """
    log = get_logger("insert_image")
    image_bytes: bytes | None = None

    try:
        # Step 1: Load Image Data into Bytes
        if isinstance(image_data, bytes):
            log.info("Loading image from in-memory bytes.")
            image_bytes = image_data
        elif isinstance(image_data, str):
            if image_data.startswith("data:image/") or re.fullmatch(
                r"^[A-Za-z0-9+/=]+$", image_data
            ):
                log.info("Attempting to decode image from Base64 string.")
                image_bytes = _decode_base64_to_bytes(image_data)

            if image_bytes is None:
                image_path_or_uri = image_data
                if image_path_or_uri.startswith("gs://"):
                    log.info(f"Downloading image from GCS: {image_path_or_uri}")
                    client = get_gcs_client()
                    if not client:
                        raise RuntimeError("GCS client not available.")
                    bucket_name, blob_name = image_path_or_uri[5:].split("/", 1)
                    blob = client.bucket(bucket_name).blob(blob_name)
                    image_bytes = blob.download_as_bytes()
                elif os.path.exists(image_path_or_uri):
                    log.info(f"Reading local image: {image_path_or_uri}")
                    with open(image_path_or_uri, "rb") as f:
                        image_bytes = f.read()
                else:
                    log.error(
                        f"Image path/URI not found or invalid: {image_path_or_uri}"
                    )
                    return
        else:
            log.error(
                f"Invalid image_data type: {type(image_data)}. Must be str or bytes."
            )
            return

        if not image_bytes:
            log.error("Image data is empty after processing.")
            return

        image_stream = io.BytesIO(image_bytes)

        # Step 2: Calculate Dimensions and Position
        # CASE A: A specific box (placeholder) is provided.
        if box_hint:
            log.info(
                "Fitting image into provided box while preserving aspect ratio."
            )

            box_left, box_top, box_width, box_height = box_hint

            # Get the original image's dimensions
            with Image.open(image_stream) as im:
                img_width_px, img_height_px = im.size

            img_width_emu = Inches(img_width_px / 96.0)
            img_height_emu = Inches(img_height_px / 96.0)

            # Calculate the scaling factor to fit without stretching
            scale = min(
                box_width / max(img_width_emu, 1),
                box_height / max(img_height_emu, 1),
            )

            final_width = int(img_width_emu * scale)
            final_height = int(img_height_emu * scale)

            # Calculate position to center the image within the box
            final_left = int(box_left + (box_width - final_width) / 2)
            final_top = int(box_top + (box_height - final_height) / 2)

            image_stream.seek(0)
            slide.shapes.add_picture(
                image_stream,
                final_left,
                final_top,
                width=final_width,
                height=final_height,
            )

        # CASE B (Fallback): No box is provided, so center on the slide.
        else:
            log.warning("No box_hint provided. Using auto-centering logic.")
            with Image.open(image_stream) as im:
                iw_px, ih_px = im.size

            slide_width = prs.slide_width
            slide_height = prs.slide_height
            iw_emu = Inches(iw_px / 96.0)
            ih_emu = Inches(ih_px / 96.0)

            scale = min(
                (slide_width * 0.7) / max(iw_emu, 1),
                (slide_height * 0.7) / max(ih_emu, 1),
            )
            final_width = int(iw_emu * scale)
            final_height = int(ih_emu * scale)
            left = int((slide_width - final_width) // 2)
            top = int((slide_height - final_height) // 2)

            if final_width > 0 and final_height > 0:
                image_stream.seek(0)
                slide.shapes.add_picture(
                    image_stream,
                    left,
                    top,
                    width=final_width,
                    height=final_height,
                )
            else:
                log.warning(
                    "Calculated image dimensions are non-positive. Skipping."
                )

        log.info("Successfully inserted image onto slide.")
    except Exception as e:
        log.error(f"Failed to insert image: {e}", exc_info=True)
