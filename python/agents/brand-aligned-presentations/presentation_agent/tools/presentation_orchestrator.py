# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import asyncio
import json
import os
import tempfile
import uuid

from google.adk.tools.tool_context import ToolContext
from google.genai import types
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches

from ..shared_libraries.config import (
    DEFAULT_TEMPLATE_URI,
    GCS_BUCKET_NAME,
    get_logger,
)
from ..shared_libraries.models import CoverSpec, DeckSpec, SlideSpec
from ..shared_libraries.utils import _insert_image
from .artifact_utils import get_gcs_file_as_local_path, save_presentation
from .pptx_editor import _insert_visual_into_slide
from .visual_generator import generate_visual


def get_smart_layout(prs: Presentation, requested_name: str):
    """
    Intelligently maps a requested layout name to the best available layout in the current presentation.
    """
    log = get_logger("layout_mapper")
    requested_name = requested_name.lower()

    #  Avoid "Title and Chart" as it requies data (demo ignore this) placeholders
    if "chart" in requested_name:
        log.info(
            f"Overriding layout request '{requested_name}' to 'Title and Image' for better spacing."
        )
        requested_name = "title and image"

    layouts = prs.slide_layouts

    # 1. Exact Match
    for layout in layouts:
        if layout.name.lower() == requested_name:
            return layout

    # 2. Key-word Based Mapping
    mapping = [
        ("two", ["two content", "comparison", "side by side", "dual", "split"]),
        (
            "image",
            [
                "image",
                "picture",
                "visual",
                "photo",
                "graphic",
                "title and image",
                "image grid",
            ],
        ),
        ("quote", ["quote", "testimonial", "statement"]),
        ("agenda", ["agenda", "toc", "roadmap", "contents"]),
        ("section", ["section header", "divider", "transition"]),
        (
            "content",
            [
                "title and content",
                "content slide",
                "standard body",
                "bullet",
                "subhead",
                "left",
            ],
        ),
        ("closing", ["closing", "thank you", "end", "contact"]),
        ("title", ["title slide", "cover", "intro", "opening", "only"]),
    ]

    # Check for conceptual matches
    for concept, keywords in mapping:
        if concept in requested_name:
            for layout in layouts:
                if any(k in layout.name.lower() for k in keywords):
                    log.info(
                        f"Mapped '{requested_name}' to layout concept '{concept}' -> '{layout.name}'"
                    )
                    return layout

    # 3. Structural Fallbacks based on common PowerPoint ordering
    try:
        # If it's clearly a title slide
        if "title" in requested_name or "cover" in requested_name:
            # Look for index 0 (Standard Title) or anything with 'title'
            for layout in layouts:
                if (
                    "title" in layout.name.lower()
                    and "content" not in layout.name.lower()
                ):
                    return layout
            return layouts[0]

        # Default to standard content layout (usually index 1)
        # However, we MUST avoid "Title Only" or "Blank" if we expect content
        final_layout = layouts[1] if len(layouts) > 1 else layouts[0]

        if any(
            k in requested_name for k in ["content", "chart", "image", "bullet"]
        ):
            if (
                "only" in final_layout.name.lower()
                or "blank" in final_layout.name.lower()
            ):
                # Try to find a layout that actually has content placeholders
                for layout in layouts:
                    if (
                        "content" in layout.name.lower()
                        or "body" in layout.name.lower()
                    ):
                        return layout

        return final_layout
    except IndexError:
        return layouts[0]


async def render_deck_from_spec(
    spec_dict: dict,
    out_pptx: str,
    tool_context: ToolContext,
    template_pptx: str | None = None,
) -> str:
    """
    Renders a presentation from a spec, using a template if provided or falling
    back to a default style. Saves to a temp file and returns the path.
    """
    log = get_logger("render_deck_from_spec")
    try:
        # 1. WORKING TEMPLATE SELECTION
        if template_pptx and os.path.exists(template_pptx):
            log.info(
                f"Using user template '{template_pptx}' as the foundation."
            )
            working_template = template_pptx
        else:
            log.error("No valid user template provided. Aborting.")
            return "Error: No valid template provided."

        # Load the user's template directly
        prs = Presentation(working_template)

        def render_slide_content(slide, spec_obj, is_cover=False):
            """Intelligently populates a slide based on its layout and the spec's content."""

            def _rm_md(t: str) -> str:
                return t.replace("**", "")

            def _apply_bullets(tf, bullets):
                tf.clear()
                tf.vertical_anchor = MSO_ANCHOR.TOP
                for i, bullet_text in enumerate(bullets):
                    # Check for sub-bullets natively
                    is_sub_bullet = (
                        bullet_text.startswith("  ")
                        or bullet_text.startswith("\t")
                        or bullet_text.startswith("- ")
                    )

                    # Clean up the bullet prefix so PowerPoint handles the symbol natively
                    clean_text = bullet_text.strip(" \t-•*")

                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.level = 1 if is_sub_bullet else 0

                    # Parse **bold** to make well-organized headers
                    parts = clean_text.split("**")
                    for j, part in enumerate(parts):
                        if not part:
                            continue
                        r = p.add_run()
                        r.text = part
                        if j % 2 != 0:
                            r.font.bold = True

            # TITLE & SUBTITLE
            # Robustly find placeholders directly from the slide's shapes
            title_ph = None
            subhead_ph = None
            content_phs = []

            for shape in slide.placeholders:
                if shape.placeholder_format.type in (
                    PP_PLACEHOLDER_TYPE.TITLE,
                    PP_PLACEHOLDER_TYPE.CENTER_TITLE,
                ):
                    title_ph = shape
                elif (
                    shape.placeholder_format.type
                    == PP_PLACEHOLDER_TYPE.SUBTITLE
                ):
                    subhead_ph = shape
                elif shape.placeholder_format.type in (
                    PP_PLACEHOLDER_TYPE.BODY,
                    PP_PLACEHOLDER_TYPE.OBJECT,
                    PP_PLACEHOLDER_TYPE.PICTURE,
                ):
                    content_phs.append(shape)

            if not title_ph and slide.shapes.title:
                title_ph = slide.shapes.title

            if title_ph and hasattr(spec_obj, "title") and spec_obj.title:
                # Use standard text assignment to preserve inheritance
                title_ph.text = _rm_md(spec_obj.title)
                if is_cover and title_ph.text_frame.paragraphs:
                    title_ph.text_frame.paragraphs[
                        0
                    ].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

            if subhead_ph and hasattr(spec_obj, "subhead") and spec_obj.subhead:
                subhead_ph.text = _rm_md(spec_obj.subhead)

            if is_cover:
                return

            # CONTENT & VISUALS
            has_bullets = hasattr(spec_obj, "bullets") and bool(
                spec_obj.bullets
            )
            image_source = getattr(spec_obj, "image_data", None) or getattr(
                spec_obj, "image_file_path", None
            )

            # Find all potential content placeholders, sorted left-to-right
            phs = sorted(content_phs, key=lambda p: p.left)

            # 1. We have BOTH Text and an Image to render
            if has_bullets and image_source:
                if len(phs) >= 2:
                    # Ideal Scenario: The user's layout actually has two placeholders.
                    text_ph = phs[0]
                    img_ph = phs[1]

                    # Write Text
                    _apply_bullets(text_ph.text_frame, spec_obj.bullets)

                    # Insert Image
                    _insert_visual_into_slide(
                        prs, slide, image_source, target_placeholder=img_ph
                    )

                elif len(phs) == 1:
                    # HACK SCENARIO: Resize 1-col layout
                    text_ph = phs[0]

                    # Squeeze the text box to the left
                    original_width = text_ph.width
                    text_ph.width = int(original_width * 0.45)

                    # Write Text
                    _apply_bullets(text_ph.text_frame, spec_obj.bullets)

                    # Float the image on the right with safer margins
                    box_hint = (
                        text_ph.left + text_ph.width + int(Inches(0.3)),
                        text_ph.top + int(Inches(0.2)),
                        int(original_width * 0.45),
                        text_ph.height - int(Inches(0.2)),
                    )
                    _insert_image(prs, slide, image_source, box_hint=box_hint)
                    log.info(
                        "Dynamically adjusted 1-column layout to fit both text and image safely."
                    )

            # 2. We ONLY have Text
            elif has_bullets and not image_source:
                if len(phs) >= 1:
                    text_ph = phs[0]

                    # Write Text
                    _apply_bullets(text_ph.text_frame, spec_obj.bullets)

            # 3. We ONLY have an Image
            elif not has_bullets and image_source:
                if len(phs) >= 1:
                    _insert_visual_into_slide(
                        prs, slide, image_source, target_placeholder=phs[0]
                    )
                else:
                    # Float it if no placeholder exists at all
                    box_hint = (
                        int(prs.slide_width * 0.1),
                        int(prs.slide_height * 0.2),
                        int(prs.slide_width * 0.8),
                        int(prs.slide_height * 0.7),
                    )
                    _insert_image(prs, slide, image_source, box_hint=box_hint)

        # 2. GENERATE SLIDES
        # We handle the Cover separately. If the template already has a slide,
        # we treat it as the Cover Page and simply update its text.
        # If it's empty, we create a new Cover Slide.

        # Cover
        cover_data = spec_dict.get(
            "cover", {"title": "Strategic Research & Analysis"}
        )
        cover_spec = CoverSpec(**cover_data)

        # Remove all existing slides except the cover (or remove all if we are building a new deck from scratch)
        # If building a new deck from a template, we usually want to start fresh.
        if len(prs.slides) > 0:
            # Save the first slide as cover, delete the rest.
            for i in range(len(prs.slides) - 1, 0, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]

        if len(prs.slides) > 0:
            log.info(
                "Template has an existing slide. Using it as the Cover Page."
            )
            cover_slide = prs.slides[0]
            # No need to add a slide, just render content onto existing one
            try:
                render_slide_content(cover_slide, cover_spec, is_cover=True)
            except (KeyError, Exception) as e:
                log.warning(
                    f"Could not render cover text on existing slide: {e}"
                )
        else:
            log.info("Template is empty. Generating a new Cover Page.")
            try:
                cover_slide = prs.slides.add_slide(
                    get_smart_layout(prs, "Title Slide")
                )
                render_slide_content(cover_slide, cover_spec, is_cover=True)
            except (KeyError, Exception) as e:
                log.warning(f"Could not generate/render cover slide: {e}")

        # Body Slides
        for s_data in spec_dict.get("slides", []):
            if "title" not in s_data or not s_data["title"]:
                s_data["title"] = "Slide Content"

            try:
                s_spec = SlideSpec(**s_data)
                slide = prs.slides.add_slide(
                    get_smart_layout(prs, s_spec.layout_name)
                )
                render_slide_content(slide, s_spec)
            except (KeyError, Exception) as e:
                log.error(
                    f"Failed to render slide '{s_data.get('title')}': {e}"
                )
                continue

            # Speaker Notes/Citations
            if s_spec.speaker_notes or s_spec.citations:
                try:
                    notes = slide.notes_slide.notes_text_frame
                    text = s_spec.speaker_notes or ""
                    if s_spec.citations:
                        if text:
                            text += "\n\n---\nCitations:\n"
                        else:
                            text += "Citations:\n"
                        for citation in s_spec.citations:
                            text += f"- {citation}\n"
                    notes.text = text
                except Exception:
                    pass

        # Closing
        try:
            closing_layout = get_smart_layout(prs, "Closing Slide")
            closing = prs.slides.add_slide(closing_layout)
            t_ph = closing.shapes.title or next(
                (
                    s
                    for s in closing.placeholders
                    if s.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE
                ),
                None,
            )
            if t_ph:
                t_ph.text = spec_dict.get("closing_title", "Thank You")
        except (KeyError, Exception) as e:
            log.warning(f"Could not generate closing slide: {e}")

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            return tmp.name
    except Exception as e:
        log.error(f"Render failed: {e}", exc_info=True)
        return f"Error: Render failed. {e}"


async def generate_and_render_deck(
    tool_context: ToolContext,
    deck_spec: dict | None = None,
    spec_artifact_name: str | None = None,
    template_path: str | None = None,
) -> dict:
    """
    Orchestrates the entire deck generation process.
    """
    log = get_logger("generate_and_render_deck_tool")
    try:
        spec_dict = deck_spec

        # 1. PRIORITY: Load from Session State (Invisible in UI)
        if not spec_dict and not spec_artifact_name:
            spec_dict = tool_context.state.get("current_deck_spec")
            if spec_dict:
                log.info("Loaded DeckSpec from session state.")

        # 2. FALLBACK: Load from Artifact Store (Explicitly named file - for manually uploaded plans)
        if not spec_dict and spec_artifact_name:
            log.info(f"Loading DeckSpec from artifact: '{spec_artifact_name}'")
            try:
                artifact = await tool_context.load_artifact(spec_artifact_name)

                # Propagation Retry for named artifacts too
                if not artifact:
                    log.warning(
                        f"Artifact '{spec_artifact_name}' not found. Waiting 2s..."
                    )
                    await asyncio.sleep(2.0)
                    artifact = await tool_context.load_artifact(
                        spec_artifact_name
                    )

                if artifact:
                    spec_json = (
                        artifact.inline_data.data
                        if isinstance(artifact, types.Part)
                        else artifact
                    )
                    if isinstance(spec_json, (bytes, bytearray)):
                        spec_dict = json.loads(spec_json.decode("utf-8"))
                    elif isinstance(spec_json, str):
                        spec_dict = json.loads(spec_json)
                    else:
                        spec_dict = spec_json
            except Exception as e:
                log.error(f"Failed to load named spec artifact: {e}")

        if not spec_dict:
            return {
                "status": "Failed",
                "message": "No active presentation plan found in session state. "
                "Please provide deck_spec or ensure an outline was generated.",
            }

        # 3. GCS-FALLBACK TEMPLATE RECOVERY
        working_template = template_path
        if not working_template or not os.path.exists(working_template):
            log.info(
                "Template path invalid or lost. Re-downloading from GCS..."
            )
            working_template = await get_gcs_file_as_local_path(
                DEFAULT_TEMPLATE_URI
            )

        # Standardize structure
        if isinstance(spec_dict.get("slides"), dict):
            spec_dict["slides"] = list(spec_dict["slides"].values())
        if "closing_title" not in spec_dict:
            spec_dict["closing_title"] = "Thank You"

        validated_spec = DeckSpec(**spec_dict)

        all_content = [validated_spec.cover] + validated_spec.slides

        # Allow up to 5 visuals per presentation
        hard_limit = 5

        visuals_kept = 0
        for slide in validated_spec.slides:
            if slide.visual_prompt:
                if visuals_kept < hard_limit:
                    visuals_kept += 1
                    # REMOVED "Title and Chart" from allowed list to prevent squeezed content
                    if slide.layout_name not in [
                        "Title and Image",
                        "Two Content",
                        "Comparison",
                    ]:
                        slide.layout_name = "Title and Image"
                else:
                    # Strip excess visuals programmatically
                    slide.visual_prompt = None

        tasks = []
        slides = []
        for item in all_content:
            if hasattr(item, "visual_prompt") and item.visual_prompt:
                tasks.append(
                    asyncio.create_task(
                        asyncio.wait_for(
                            generate_visual(item.visual_prompt), timeout=60.0
                        )
                    )
                )
                slides.append(item)

        images = await asyncio.gather(*tasks, return_exceptions=True)
        for s, img in zip(slides, images):
            if not isinstance(img, Exception):
                s.image_data = img

        out_name = f"{validated_spec.cover.title}_{uuid.uuid4().hex[:6]}.pptx"

        local_path = await render_deck_from_spec(
            validated_spec.model_dump(),
            out_name,
            tool_context,
            working_template,
        )
        if local_path.startswith("Error:"):
            return {"status": "Failed", "message": local_path}

        msg = await save_presentation(
            tool_context, out_name, local_path, GCS_BUCKET_NAME
        )
        os.remove(local_path)
        return {"status": "Success", "message": msg}
    except Exception as e:
        log.error(f"Generation failed: {e}", exc_info=True)
        return {"status": "Failed", "message": str(e)}
