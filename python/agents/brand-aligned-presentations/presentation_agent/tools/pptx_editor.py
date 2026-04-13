# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import os
from typing import Any, Literal

from google.adk.tools.tool_context import ToolContext
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.util import Inches

# Local Application Imports
from ..shared_libraries.config import get_logger
from ..shared_libraries.models import StyleProfile
from ..shared_libraries.utils import _insert_image


def edit_slide_text(
    pptx_path: str,
    slide_number: int,
    new_title: str | None = None,
    new_bullets: list[str] | None = None,
    new_speaker_notes: str | None = None,
) -> str:
    """Edits the text and speaker notes of a specific slide. Modifies the file IN-PLACE."""
    if not os.path.exists(pptx_path):
        return f"Error: File not found at {pptx_path}"
    try:
        prs = Presentation(pptx_path)
        if not (1 <= slide_number <= len(prs.slides)):
            return f"Error: Invalid slide number. Must be between 1 and {len(prs.slides)}."
        slide = prs.slides[slide_number - 1]

        try:
            if new_title is not None and slide.shapes.title:
                slide.shapes.title.text_frame.text = new_title
        except (KeyError, AttributeError):
            pass

        if new_bullets is not None:
            try:
                body_ph = next(
                    (
                        s
                        for s in slide.placeholders
                        if s.placeholder_format.type
                        in (
                            PP_PLACEHOLDER_TYPE.BODY,
                            PP_PLACEHOLDER_TYPE.OBJECT,
                        )
                    ),
                    None,
                )
                if body_ph:
                    tf = body_ph.text_frame
                    tf.clear()
                    if new_bullets:
                        tf.text = new_bullets[0]
                        for text in new_bullets[1:]:
                            tf.add_paragraph().text = text
            except (KeyError, AttributeError):
                pass

        if new_speaker_notes is not None:
            try:
                slide.notes_slide.notes_text_frame.text = new_speaker_notes
            except Exception:
                pass

        prs.save(pptx_path)
        return f"Successfully edited slide {slide_number}."
    except Exception as e:
        return f"Error: Could not edit slide text. {e}"


def add_slide_to_end(
    pptx_path: str,
    title: str,
    layout_name: str | None = None,
    bullets: list[str] | None = None,
    speaker_notes: str | None = None,
) -> str:
    """Adds a new slide to the END of the presentation. Modifies file IN-PLACE."""
    if not os.path.exists(pptx_path):
        return f"Error: File not found at {pptx_path}"
    try:
        prs = Presentation(pptx_path)
        layout_index = 1 if bullets else 5
        # Ensure we don't crash if layouts are missing
        if layout_index >= len(prs.slide_layouts):
            layout_index = 0
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        try:
            if slide.shapes.title:
                slide.shapes.title.text_frame.text = title
        except (KeyError, AttributeError):
            pass

        if bullets:
            try:
                body_ph = next(
                    (
                        s
                        for s in slide.placeholders
                        if s.placeholder_format.type
                        in (
                            PP_PLACEHOLDER_TYPE.BODY,
                            PP_PLACEHOLDER_TYPE.OBJECT,
                        )
                    ),
                    None,
                )
                if body_ph:
                    tf = body_ph.text_frame
                    tf.clear()
                    tf.text = bullets[0]
                    for text in bullets[1:]:
                        tf.add_paragraph().text = text
            except (KeyError, AttributeError):
                pass

        if speaker_notes:
            try:
                slide.notes_slide.notes_text_frame.text = speaker_notes
            except Exception:
                pass

        prs.save(pptx_path)
        return f"Successfully added new slide '{title}'."
    except Exception as e:
        return f"Error: Could not add new slide. {e}"


def delete_slide(pptx_path: str, slide_number_to_delete: int) -> str:
    """Deletes a specific slide. Modifies the file IN-PLACE."""
    if not os.path.exists(pptx_path):
        return f"Error: File not found at {pptx_path}"
    try:
        prs = Presentation(pptx_path)
        if not (1 <= slide_number_to_delete <= len(prs.slides)):
            return f"Error: Invalid slide number. Must be between 1 and {len(prs.slides)}."
        rId = prs.slides._sldIdLst[slide_number_to_delete - 1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[slide_number_to_delete - 1]
        prs.save(pptx_path)
        return f"Successfully deleted slide {slide_number_to_delete}."
    except Exception as e:
        return f"Error: Could not delete slide. {e}"


def replace_slide_visual(
    pptx_path: str,
    slide_number: int,
    image_data_base64: str,
    target_type: Literal["picture", "chart", "any"] | None = "any",
) -> str:
    """
    Replaces an existing visual element (Picture or Chart) on a slide with a new image.
    If 'chart' is specified, it targets charts first and preserves their exact dimensions.
    If 'picture' is specified, it targets existing pictures.
    If 'any' is specified (default), it looks for charts first, then pictures.
    If no matching visual is found, it inserts a new centered image.
    """
    if not os.path.exists(pptx_path):
        return f"Error: File not found at {pptx_path}"
    try:
        prs = Presentation(pptx_path)
        if not (1 <= slide_number <= len(prs.slides)):
            return f"Error: Invalid slide number {slide_number}."
        slide = prs.slides[slide_number - 1]

        target_shape = None
        box_hint = None

        # 1. Search for a Chart if allowed
        if target_type in ("chart", "any"):
            target_shape = next(
                (
                    s
                    for s in slide.shapes
                    if s.shape_type == MSO_SHAPE_TYPE.CHART
                ),
                None,
            )

        # 2. Search for a Picture if no chart was found or if explicitly requested
        if not target_shape and target_type in ("picture", "any"):
            target_shape = next(
                (
                    s
                    for s in slide.shapes
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE
                ),
                None,
            )

        # 3. If a shape was found, capture its dimensions and remove it
        if target_shape:
            box_hint = (
                target_shape.left,
                target_shape.top,
                target_shape.width,
                target_shape.height,
            )
            sp = target_shape._element
            sp.getparent().remove(sp)
            # Remove other pictures if 'picture' was explicitly requested to avoid overlap
            if target_type == "picture":
                for other_pic in [
                    s
                    for s in slide.shapes
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE
                ]:
                    sp_other = other_pic._element
                    sp_other.getparent().remove(sp_other)

        # 4. Insert the new image
        _insert_image(prs, slide, image_data_base64, box_hint=box_hint)
        prs.save(pptx_path)

        type_msg = "visual" if target_type == "any" else target_type
        return f"Successfully replaced {type_msg} on slide {slide_number}."
    except Exception as e:
        return f"Error: Could not replace visual. {e}"


async def update_element_layout(
    pptx_path: str,
    slide_number: int,
    element_index: int = 0,
    left_inches: float | None = None,
    top_inches: float | None = None,
    width_inches: float | None = None,
    height_inches: float | None = None,
) -> str:
    """Updates position and size of an element on a slide."""
    log = get_logger("update_element_layout")
    if not os.path.exists(pptx_path):
        return f"Error: File not found: {pptx_path}"
    try:
        prs = Presentation(pptx_path)
        if not (1 <= slide_number <= len(prs.slides)):
            return "Error: Invalid slide number."
        slide = prs.slides[slide_number - 1]
        target_shapes = [s for s in slide.shapes if not s.is_placeholder]
        if element_index >= len(target_shapes):
            return f"Error: Element at index {element_index} not found on slide {slide_number}."
        shape = target_shapes[element_index]
        if left_inches is not None:
            shape.left = Inches(left_inches)
        if top_inches is not None:
            shape.top = Inches(top_inches)
        if width_inches is not None:
            shape.width = Inches(width_inches)
        if height_inches is not None:
            shape.height = Inches(height_inches)
        prs.save(pptx_path)
        return f"Updated layout for element {element_index} on slide {slide_number}."
    except Exception as e:
        log.error(f"Failed to update element layout: {e}", exc_info=True)
        return f"Error: Could not update layout. {e}"


async def inspect_template_style(
    tool_context: ToolContext, template_pptx: str
) -> dict:
    """
    Extracts style information (fonts, colors) from a template's slide master.
    Can accept a local path or an artifact name.
    """
    log = get_logger("style_tool")

    # Sensible defaults
    defaults = {
        "title_font_name": "Calibri Light",
        "title_font_size_pt": 28.0,
        "title_font_color_rgb": (0, 0, 0),
        "body_font_name": "Calibri",
        "body_font_size_pt": 18.0,
        "body_font_color_rgb": (40, 40, 40),
    }

    resolved_path = template_pptx

    try:
        # 1. Resolve if it's not a local file
        if not os.path.exists(template_pptx):
            log.info(
                f"Path '{template_pptx}' not found locally. Attempting artifact resolution..."
            )
            from .artifact_utils import get_artifact_as_local_path

            res = await get_artifact_as_local_path(tool_context, template_pptx)
            if not res.startswith("Error"):
                resolved_path = res
            else:
                log.warning(
                    f"Could not resolve '{template_pptx}' as artifact. Using defaults."
                )
                defaults["supports_subtitles"] = False
                return defaults

        # 2. Inspect the file
        prs = Presentation(resolved_path)
        master = prs.slide_masters[0]

        def get_font_style(placeholder_type, type_key):
            try:
                for ph in master.placeholders:
                    if ph.placeholder_format.type == placeholder_type:
                        try:
                            font = ph.text_frame.paragraphs[0].font
                            name = (
                                font.name or defaults[f"{type_key}_font_name"]
                            )
                            size = (
                                font.size.pt
                                if font.size
                                else defaults[f"{type_key}_font_size_pt"]
                            )
                            # Handle color being None or not having rgb
                            color = (
                                font.color.rgb
                                if font.color
                                and hasattr(font.color, "rgb")
                                and font.color.rgb is not None
                                else defaults[f"{type_key}_font_color_rgb"]
                            )
                            return name, size, color
                        except Exception:
                            log.warning(
                                f"Could not extract {type_key} style from placeholder."
                            )
                            break
            except (KeyError, AttributeError):
                pass
            return (
                defaults[f"{type_key}_font_name"],
                defaults[f"{type_key}_font_size_pt"],
                defaults[f"{type_key}_font_color_rgb"],
            )

        t_name, t_size, t_color = get_font_style(
            PP_PLACEHOLDER_TYPE.TITLE, "title"
        )
        b_name, b_size, b_color = get_font_style(
            PP_PLACEHOLDER_TYPE.BODY, "body"
        )

        # Check if the template supports subtitles
        supports_subtitles = False
        for layout in prs.slide_layouts:
            try:
                for shape in layout.placeholders:
                    if (
                        shape.placeholder_format.type
                        == PP_PLACEHOLDER_TYPE.SUBTITLE
                    ):
                        supports_subtitles = True
                        break
            except (KeyError, AttributeError):
                continue
            if supports_subtitles:
                break

        profile = StyleProfile(
            title_font_name=t_name,
            title_font_size_pt=t_size,
            title_font_color_rgb=t_color,
            body_font_name=b_name,
            body_font_size_pt=b_size,
            body_font_color_rgb=b_color,
            accent_colors=[],
            image_box_hint=None,
            supports_subtitles=supports_subtitles,
        )
        log.info(f"Extracted style profile: {profile.dict()}")
        return profile.dict()
    except Exception as e:
        log.error(f"Style extraction failed: {e}", exc_info=True)
        defaults["supports_subtitles"] = False
        return defaults


def _insert_visual_into_slide(
    prs: Presentation,
    slide: Any,
    image_source: str | bytes,
    target_placeholder: Any | None = None,
):
    """Inserts a visual into a slide using a placeholder or centering it."""
    get_logger("insert_visual_into_slide")

    # 1. Select the placeholder to use
    placeholder_to_use = target_placeholder
    if not placeholder_to_use:
        try:
            placeholder_to_use = next(
                (
                    p
                    for p in slide.placeholders
                    if p.placeholder_format.type
                    in (
                        PP_PLACEHOLDER_TYPE.PICTURE,
                        PP_PLACEHOLDER_TYPE.OBJECT,
                        PP_PLACEHOLDER_TYPE.BODY,
                    )
                ),
                None,
            )
        except (KeyError, AttributeError):
            placeholder_to_use = None

    # 2. If a placeholder is found, use its dimensions to call _insert_image
    if placeholder_to_use:
        try:
            if placeholder_to_use.has_text_frame:
                placeholder_to_use.text_frame.clear()
            box_hint = (
                placeholder_to_use.left,
                placeholder_to_use.top,
                placeholder_to_use.width,
                placeholder_to_use.height,
            )
            _insert_image(prs, slide, image_source, box_hint=box_hint)
        except (KeyError, AttributeError, Exception):
            _insert_image(prs, slide, image_source, box_hint=None)
    else:
        # Fallback if no placeholder was ever found
        _insert_image(prs, slide, image_source, box_hint=None)


def get_safe_layout(prs, requested_layout_name):
    """Safely retrieves a layout using keyword-based mapping."""
    if not requested_layout_name:
        return prs.slide_layouts[0]

    requested_name = requested_layout_name.lower()
    layouts = prs.slide_layouts

    # 1. Exact Match
    for layout in layouts:
        if layout.name.lower() == requested_name:
            return layout

    # 2. Key-word Based Mapping
    mapping = [
        ("two", ["two content", "comparison", "side by side", "dual", "split"]),
        ("chart", ["chart", "graph", "plot", "data", "statistic"]),
        ("image", ["image", "picture", "visual", "photo", "graphic"]),
        (
            "content",
            ["title and content", "content slide", "body", "bullet", "subhead"],
        ),
        ("closing", ["closing", "thank you", "end", "contact"]),
        ("title", ["title slide", "cover", "intro", "opening", "only"]),
    ]

    for concept, keywords in mapping:
        if concept in requested_name:
            for layout in layouts:
                if any(k in layout.name.lower() for k in keywords):
                    return layout

    # 3. Fallbacks
    try:
        if "title" in requested_name or "cover" in requested_name:
            for layout in layouts:
                if (
                    "title" in layout.name.lower()
                    and "content" not in layout.name.lower()
                ):
                    return layout
            return layouts[0]
        return layouts[1] if len(layouts) > 1 else layouts[0]
    except Exception:
        return layouts[0]


def read_presentation_outline(pptx_path: str) -> str:
    """Reads the slide titles from a presentation to provide a 'table of contents'."""
    import os

    from pptx import Presentation

    if not os.path.exists(pptx_path):
        return f"Error: File not found at {pptx_path}"
    try:
        prs = Presentation(pptx_path)
        titles = []
        for i, slide in enumerate(prs.slides):
            title_text = f"Slide {i + 1}: [Untitled]"
            try:
                if (
                    slide.shapes.title
                    and slide.shapes.title.has_text_frame
                    and slide.shapes.title.text_frame.text
                ):
                    title_text = (
                        f"Slide {i + 1}: {slide.shapes.title.text_frame.text}"
                    )
            except (KeyError, AttributeError):
                pass
            titles.append(title_text)
        return "\n".join(titles) if titles else "The presentation is empty."
    except Exception as e:
        return f"Error: Could not read presentation outline. {e}"


def read_presentation_details(pptx_path: str) -> str:
    """Extracts slide titles and all text content from a .pptx file."""
    import os

    from pptx import Presentation

    if not os.path.exists(pptx_path):
        return f"Error: File not found at {pptx_path}"
    try:
        prs = Presentation(pptx_path)
        if not prs.slides:
            return "The presentation is empty."
        details = []
        for i, slide in enumerate(prs.slides):
            title = "[Untitled]"
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title = slide.shapes.title.text_frame.text
            except (KeyError, AttributeError):
                pass

            content_lines = []
            for s in slide.shapes:
                try:
                    if s.has_text_frame:
                        # Skip title as we already have it
                        is_title = False
                        try:
                            if s == slide.shapes.title:
                                is_title = True
                        except (KeyError, AttributeError):
                            pass

                        if not is_title:
                            for p in s.text_frame.paragraphs:
                                if p.runs:
                                    content_lines.append(
                                        f"- {''.join(r.text for r in p.runs).strip()}"
                                    )
                except (KeyError, AttributeError):
                    continue

            content = "\n".join(content_lines)
            details.append(
                f"Slide {i + 1} Title: {title}\nContent:\n{content or '[No text content found]'}"
            )
        return "\n\n---\n\n".join(details)
    except Exception as e:
        return f"Error: Could not read presentation details. {e}"


def extract_slide_content(file_path: str, slide_number: int) -> str:
    """
    Extracts all text content and speaker notes from a specific slide in a PowerPoint presentation.
    """
    import os

    from pptx import Presentation

    if not os.path.exists(file_path):
        return f"Error: The file at '{file_path}' was not found."

    try:
        presentation = Presentation(file_path)

        num_slides = len(presentation.slides)
        if not 1 <= slide_number <= num_slides:
            return (
                f"Invalid slide number: {slide_number}. "
                f"Presentation has {num_slides} slides (1 to {num_slides})."
            )

        slide_index = slide_number - 1
        slide = presentation.slides[slide_index]

        slide_text = []
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    run_texts = [run.text for run in paragraph.runs]
                    slide_text.append("".join(run_texts))
            except (KeyError, AttributeError):
                continue

        # Extract speaker notes
        speaker_notes = ""
        try:
            if slide.has_notes_slide:
                speaker_notes = slide.notes_slide.notes_text_frame.text
        except Exception:
            pass

        content = "\n".join(slide_text)
        if speaker_notes:
            content += f"\n\nSPEAKER NOTES:\n{speaker_notes}"

        return content

    except Exception as e:
        return f"Error: Could not extract slide content. Details: {e}"


def extract_images_from_slide(file_path: str, slide_number: int) -> list:
    """
    Extracts all images from a specific slide in a PowerPoint presentation.
    """
    import os

    from pptx import Presentation

    if not os.path.exists(file_path):
        return [f"Error: The file at '{file_path}' was not found."]

    try:
        presentation = Presentation(file_path)

        num_slides = len(presentation.slides)
        if not 1 <= slide_number <= num_slides:
            return [f"Invalid slide number: {slide_number}."]

        slide_index = slide_number - 1
        slide = presentation.slides[slide_index]

        images = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "image"):
                    image_blob = shape.image.blob
                    content_type = shape.image.content_type

                    images.append(
                        {"blob": image_blob, "content_type": content_type}
                    )
            except (KeyError, AttributeError):
                continue

        return images

    except Exception as e:
        return [f"Error: {e}"]
