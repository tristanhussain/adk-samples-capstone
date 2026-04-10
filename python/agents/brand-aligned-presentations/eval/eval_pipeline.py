# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import asyncio
import json
import os
import re
from io import BytesIO

from dotenv import load_dotenv
from google import genai
from google.genai import types
from pptx import Presentation

# Load env before importing app to ensure we can manipulate it if needed
load_dotenv()

# Force InMemory services for evaluation
os.environ["GCP_STAGING_BUCKET"] = ""
os.environ["LOCAL_DEV"] = "true"

from presentation_agent.agent import PresentationExpertApp


def create_valid_mock_pptx():
    """Generates bytes for a valid 5-slide PowerPoint file."""
    prs = Presentation()
    for i in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = f"Slide {i} Title"
        slide.placeholders[
            1
        ].text = f"This is the body content for slide {i}. It contains useful data for testing."

    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()


async def mock_upload_artifact(app, session_id, filename):
    """Mocks an 'upload' by saving a valid PPTX into the agent's session store."""
    content = create_valid_mock_pptx()

    artifact = types.Part(
        inline_data=types.Blob(
            data=content,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    )

    # Save the artifact using the EXACT keyword arguments from ADK signature
    await app._runner.artifact_service.save_artifact(
        app_name="presentation_agent",
        user_id="evaluator",
        filename=filename,
        artifact=artifact,
        session_id=session_id,
    )
    return filename


async def evaluate_scenario(client, app, test_case, session_id):
    """Runs a single test scenario and calculates metrics."""
    print(f"\n--- Running Scenario: {test_case['name']} ---")

    # Pre-load artifacts if the test case requires them
    if "artifacts" in test_case:
        for filename in test_case["artifacts"]:
            await mock_upload_artifact(app, session_id, filename)
            print(f"  [Setup] Mocked upload of valid PPTX: {filename}")

    actual_tool_calls = []
    final_response = ""
    generated_deck_spec = None

    # --- TURN 1: Initial Request ---
    print(f"Turn 1: Requesting... '{test_case['prompt']}'")
    message = types.Content(
        role="user", parts=[types.Part.from_text(text=test_case["prompt"])]
    )

    # Use the real application's runner
    events = app._runner.run(
        user_id="evaluator", session_id=session_id, new_message=message
    )
    for event in events:
        if event.content and event.content.parts:
            for part in event.content.parts:
                if part.function_call:
                    tool_name = part.function_call.name
                    actual_tool_calls.append(tool_name)
                    print(
                        f"  [Tool Call]: {tool_name}({part.function_call.args})"
                    )
                if part.text:
                    final_response += part.text

    # --- TURN 2: Revision (if requested) ---
    if test_case.get("requires_revision", False):
        rev_prompt = test_case["revision_prompt"]
        print(f"Turn 2 (Revision): Requesting... '{rev_prompt}'")
        rev_msg = types.Content(
            role="user", parts=[types.Part.from_text(text=rev_prompt)]
        )
        events_rev = app._runner.run(
            user_id="evaluator", session_id=session_id, new_message=rev_msg
        )
        for event in events_rev:
            if event.content and event.content.parts:
                for part in event.content.parts:
                    if part.text:
                        final_response += part.text
                    if part.function_call:
                        tool_name = part.function_call.name
                        actual_tool_calls.append(tool_name)
                        print(
                            f"  [Tool Call]: {tool_name}({part.function_call.args})"
                        )

    # --- FINAL TURN: Follow-up / Approval (if needed) ---
    if test_case.get("requires_approval", False):
        print("Final Turn: Sending 'Approve' to trigger rendering...")
        approve_msg = types.Content(
            role="user",
            parts=[types.Part.from_text(text="Approve and render now.")],
        )
        events_f = app._runner.run(
            user_id="evaluator", session_id=session_id, new_message=approve_msg
        )
        for event in events_f:
            if event.content and event.content.parts:
                for part in event.content.parts:
                    if part.text:
                        final_response += part.text
                    if part.function_call:
                        tool_name = part.function_call.name
                        actual_tool_calls.append(tool_name)
                        print(
                            f"  [Tool Call]: {tool_name}({part.function_call.args})"
                        )

        # AFTER ALL TURNS: Capture final deck spec from state for constraint check
        session = await app._runner.session_service.get_session(
            app_name="presentation_agent",
            user_id="evaluator",
            session_id=session_id,
        )
        generated_deck_spec = session.state.get("current_deck_spec")
        print(f"  [Full Response Total]: {final_response}")

    # --- Metrics Calculation ---

    # A. Tool Trajectory Score
    expected = test_case["expected_tools"]

    # We use a set for actual_tool_calls to avoid duplication penalties
    unique_actual = set(actual_tool_calls)
    intersection = set(expected).intersection(unique_actual)
    tool_trajectory_avg_score = len(intersection) / max(len(expected), 1)

    # B. Response Match Score
    judge_match_prompt = f"""
    Evaluate how well the agent's final response matches the user's intent and includes requested information.
    User Intent: "{test_case["prompt"]}"
    Target Keywords: {test_case["target_keywords"]}
    Agent Response: "{final_response}"
    
    Provide a score from 0.0 (no match) to 1.0 (perfect match).
    Return ONLY the float number, no other text.
    """
    try:
        match_resp = client.models.generate_content(
            model="gemini-2.5-flash", contents=judge_match_prompt
        )
        match = re.search(r"(\d\.\d+)", match_resp.text)
        if not match:
            match = re.search(r"\b([01])\b", match_resp.text)

        response_match_score = float(match.group(1)) if match else 0.0
    except Exception as e:
        print(f"  [Error] Judge match failed: {e}")
        response_match_score = 0.0

    # C. Constraint Compliance Score
    constraint_compliance_score = 0.0
    score = 0
    total_constraints = 0

    if "constraints" in test_case:
        constraints = test_case["constraints"]

        if "slide_count" in constraints and generated_deck_spec:
            total_constraints += 1
            if (
                len(generated_deck_spec.get("slides", []))
                == constraints["slide_count"]
            ):
                score += 1
            else:
                print(
                    f"  [Debug] Slide count mismatch. Expected {constraints['slide_count']},got {len(generated_deck_spec.get('slides', []))}"
                )

        if "required_terms" in constraints and generated_deck_spec:
            total_constraints += 1
            deck_text = json.dumps(generated_deck_spec).lower()
            if all(
                term.lower() in deck_text
                for term in constraints["required_terms"]
            ):
                score += 1
            else:
                print("  [Debug] Required terms missing from deck_spec.")

        if "has_citations" in constraints:
            total_constraints += 1
            # Check for URLs in speaker notes of deck_spec (State is the source of truth)
            has_citation = False
            if generated_deck_spec:
                # Check slides
                for slide in generated_deck_spec.get("slides", []):
                    notes = slide.get("speaker_notes", "").lower()
                    if "http" in notes:
                        has_citation = True
                        break
                # Check research summary specifically in state
                session = await app._runner.session_service.get_session(
                    app_name="presentation_agent",
                    user_id="evaluator",
                    session_id=session_id,
                )
                state_summary = session.state.get(
                    "research_summary", ""
                ).lower()
                if "http" in state_summary:
                    has_citation = True

            # Fallback check final response text
            if not has_citation:
                has_citation = bool(re.search(r"https?://", final_response))

            if has_citation:
                score += 1
            else:
                print(
                    "  [Debug] Citation check failed. No URLs found in state or response."
                )

    constraint_compliance_score = (
        (score / total_constraints) if total_constraints > 0 else 1.0
    )

    results = {
        "scenario": test_case["name"],
        "tool_trajectory_avg_score": round(tool_trajectory_avg_score, 2),
        "response_match_score": round(response_match_score, 2),
        "constraint_compliance_score": round(constraint_compliance_score, 2),
    }
    print(f"  [Results] {results}")
    return results


async def run_evaluation_with_metrics():
    """
    Standard Evaluation Pipeline testing Creation, Editing, and Research capabilities.
    """
    print("=========================================")
    print("🚀 Starting Multi-Scenario Eval Pipeline 🚀")
    print("=========================================\n")

    project_id = os.environ.get("GOOGLE_CLOUD_PROJECT")
    location = os.environ.get("GOOGLE_CLOUD_LOCATION", "us-central1")

    if not project_id:
        print("Skipping: GOOGLE_CLOUD_PROJECT not set.")
        return

    client = genai.Client(vertexai=True, project=project_id, location=location)

    # Initialize the Real App
    app = PresentationExpertApp()

    # 1. Define Test Scenarios
    test_scenarios = [
        {
            "name": "Standard Creation Workflow",
            "prompt": "Create a 2-slide presentation about why Cloud Computing is good. You must mention 'AWS' and 'Azure'. No external tools.",
            "requires_approval": True,
            "expected_tools": [
                "get_gcs_file_as_local_path",
                "inspect_template_style",
                "generate_and_save_outline",
                "batch_generate_slides",
                "generate_and_render_deck",
            ],
            "target_keywords": ["presentation", "ready"],
            "constraints": {
                "slide_count": 2,
                "required_terms": ["AWS", "Azure"],
            },
        },
        {
            "name": "Editing: Index Shifting",
            "prompt": "I have uploaded 'draft.pptx'. First, delete slide 2. Then, add a new slide at the end summarizing our Q3 goals.",
            "artifacts": ["draft.pptx"],
            "requires_approval": False,
            "strict_order": False,
            "expected_tools": [
                "get_artifact_as_local_path",
                "delete_slide",
                "read_presentation_outline",
                "add_slide_to_end",
            ],
            "target_keywords": ["deleted", "added"],
            "constraints": {},
        },
        {
            "name": "Revision Integrity & Citation Preservation",
            "prompt": "Create a 1-slide presentation about mRNA vaccines. Include research. Once the outline is shown, change the title of slide 1 to 'mRNA: The Future of Medicine'.",
            "requires_approval": True,
            "requires_revision": True,
            "revision_prompt": "Actually, change the title of slide 1 to 'mRNA: The Future of Medicine'. Keep all other research findings and URLs.",
            "expected_tools": [
                "google_research_grounded_tool",
                "generate_and_save_outline",
                "update_slide_in_spec",
                "batch_generate_slides",
            ],
            "target_keywords": ["mRNA", "Medicine"],
            "constraints": {"slide_count": 1, "has_citations": True},
        },
    ]

    # 2. Run Evaluations
    all_results = []

    for i, test_case in enumerate(test_scenarios):
        session_id = f"eval_pipeline_session_{i}"
        await app._runner.session_service.create_session(
            app_name="presentation_agent",
            user_id="evaluator",
            session_id=session_id,
        )
        try:
            res = await evaluate_scenario(client, app, test_case, session_id)
            all_results.append(res)
        except Exception as e:
            print(f"Error evaluating scenario '{test_case['name']}': {e}")

    #  aggregations ...
    if all_results:
        print("\n" + "=" * 50)
        print("📊 FINAL EVALUATION METRICS BY SCENARIO:")
        print(json.dumps(all_results, indent=4))

        avg_traj = sum(
            r["tool_trajectory_avg_score"] for r in all_results
        ) / len(all_results)
        avg_match = sum(r["response_match_score"] for r in all_results) / len(
            all_results
        )
        avg_const = sum(
            r["constraint_compliance_score"] for r in all_results
        ) / len(all_results)

        print("\n--- OVERALL AVERAGES ---")
        print(f"Tool Trajectory: {avg_traj:.2f}")
        print(f"Response Match:  {avg_match:.2f}")
        print(f"Constraint Comp: {avg_const:.2f}")

    print("=" * 50)


if __name__ == "__main__":
    asyncio.run(run_evaluation_with_metrics())
